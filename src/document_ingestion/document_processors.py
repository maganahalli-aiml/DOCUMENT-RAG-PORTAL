"""
Enhanced Document Processors for Multiple File Types
Supports: .ppt, .docx, .md, .txt, .pdf, .xlsx, .csv, SQL databases
"""

import os
import io
import pandas as pd
import sqlite3
from typing import List, Dict, Any, Optional, Union
from pathlib import Path
from abc import ABC, abstractmethod

# Document processing imports
import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation
import openpyxl
from PIL import Image
import pytesseract
import markdown
from bs4 import BeautifulSoup

# LangChain imports
from langchain.schema import Document as LangChainDocument
from langchain.text_splitter import RecursiveCharacterTextSplitter


class BaseDocumentProcessor(ABC):
    """Base class for all document processors"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap
        )
    
    @abstractmethod
    def process(self, file_path: str) -> List[LangChainDocument]:
        """Process the document and return chunks"""
        pass
    
    def _create_documents(self, texts: List[str], metadata: Dict[str, Any]) -> List[LangChainDocument]:
        """Create LangChain documents from text chunks"""
        chunks = self.text_splitter.split_text("\n".join(texts))
        documents = []
        
        for i, chunk in enumerate(chunks):
            doc_metadata = metadata.copy()
            doc_metadata.update({
                "chunk_id": i,
                "chunk_count": len(chunks)
            })
            documents.append(LangChainDocument(page_content=chunk, metadata=doc_metadata))
        
        return documents


class PowerPointProcessor(BaseDocumentProcessor):
    """Process PowerPoint presentations (.ppt, .pptx)"""
    
    def process(self, file_path: str) -> List[LangChainDocument]:
        """Extract text and images from PowerPoint files"""
        try:
            prs = Presentation(file_path)
            texts = []
            images = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                    
                    # Extract images and perform OCR
                    if shape.shape_type == 13:  # Picture type
                        try:
                            image_data = shape.image.blob
                            image = Image.open(io.BytesIO(image_data))
                            ocr_text = pytesseract.image_to_string(image)
                            if ocr_text.strip():
                                slide_text.append(f"[Image OCR]: {ocr_text.strip()}")
                                images.append({
                                    "slide": slide_num,
                                    "ocr_text": ocr_text.strip(),
                                    "image_size": image.size
                                })
                        except Exception as e:
                            print(f"Error processing image in slide {slide_num}: {e}")
                
                if slide_text:
                    texts.append(f"Slide {slide_num}:\n" + "\n".join(slide_text))
            
            metadata = {
                "source": file_path,
                "file_type": "powerpoint",
                "total_slides": len(prs.slides),
                "images_found": len(images),
                "images_metadata": images
            }
            
            return self._create_documents(texts, metadata)
            
        except Exception as e:
            raise Exception(f"Error processing PowerPoint file {file_path}: {e}")


class ExcelProcessor(BaseDocumentProcessor):
    """Process Excel and CSV files (.xlsx, .csv)"""
    
    def process(self, file_path: str) -> List[LangChainDocument]:
        """Extract data from Excel/CSV files with table structure preservation"""
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if file_ext == '.csv':
                df = pd.read_csv(file_path)
                sheets = {'Sheet1': df}
            else:
                excel_file = pd.ExcelFile(file_path)
                sheets = {sheet_name: pd.read_excel(file_path, sheet_name=sheet_name) 
                         for sheet_name in excel_file.sheet_names}
            
            texts = []
            tables_metadata = []
            
            for sheet_name, df in sheets.items():
                if df.empty:
                    continue
                
                # Convert DataFrame to structured text
                sheet_text = [f"Sheet: {sheet_name}"]
                
                # Add column headers
                headers = df.columns.tolist()
                sheet_text.append(f"Columns: {', '.join(headers)}")
                
                # Add summary statistics
                numeric_cols = df.select_dtypes(include=['number']).columns
                if len(numeric_cols) > 0:
                    sheet_text.append("Summary Statistics:")
                    for col in numeric_cols:
                        stats = df[col].describe()
                        # Convert numpy types to Python types for serialization
                        mean_val = float(stats['mean']) if pd.notna(stats['mean']) else 0.0
                        min_val = float(stats['min']) if pd.notna(stats['min']) else 0.0
                        max_val = float(stats['max']) if pd.notna(stats['max']) else 0.0
                        sheet_text.append(f"{col}: mean={mean_val:.2f}, min={min_val}, max={max_val}")
                
                # Add sample rows (first 5 and last 5)
                sheet_text.append("\nSample Data:")
                sample_rows = min(10, len(df))
                for idx, row in df.head(sample_rows).iterrows():
                    row_text = ", ".join([f"{col}: {val}" for col, val in row.items() if pd.notna(val)])
                    sheet_text.append(f"Row {idx + 1}: {row_text}")
                
                # Add searchable table content
                sheet_text.append("\nFull Table Content:")
                table_content = df.to_string(index=False, max_rows=100)
                sheet_text.append(table_content)
                
                texts.append("\n".join(sheet_text))
                
                # Store table metadata
                tables_metadata.append({
                    "sheet_name": sheet_name,
                    "rows": len(df),
                    "columns": len(df.columns),
                    "column_names": headers,
                    "data_types": {col: str(dtype) for col, dtype in df.dtypes.items()},  # Convert numpy types to strings
                    "memory_usage": int(df.memory_usage(deep=True).sum())  # Convert numpy int to Python int
                })
            
            metadata = {
                "source": file_path,
                "file_type": "spreadsheet",
                "total_sheets": len(sheets),
                "tables_metadata": tables_metadata
            }
            
            return self._create_documents(texts, metadata)
            
        except Exception as e:
            raise Exception(f"Error processing Excel/CSV file {file_path}: {e}")


class MarkdownProcessor(BaseDocumentProcessor):
    """Process Markdown files (.md)"""
    
    def process(self, file_path: str) -> List[LangChainDocument]:
        """Extract text and structure from Markdown files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                md_content = file.read()
            
            # Convert markdown to HTML then extract text
            html = markdown.markdown(md_content, extensions=['meta', 'toc', 'tables'])
            soup = BeautifulSoup(html, 'html.parser')
            
            # Extract metadata from markdown
            md_processor = markdown.Markdown(extensions=['meta'])
            md_processor.convert(md_content)
            file_metadata = getattr(md_processor, 'Meta', {})
            
            # Extract structured content
            texts = []
            headings = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
            
            current_section = ""
            section_content = []
            
            for element in soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6', 'p', 'ul', 'ol', 'table']):
                if element.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                    # Save previous section
                    if current_section and section_content:
                        texts.append(f"{current_section}\n" + "\n".join(section_content))
                    
                    # Start new section
                    current_section = element.get_text().strip()
                    section_content = []
                else:
                    # Add content to current section
                    content = element.get_text().strip()
                    if content:
                        section_content.append(content)
            
            # Add final section
            if current_section and section_content:
                texts.append(f"{current_section}\n" + "\n".join(section_content))
            
            # If no sections found, use the whole content
            if not texts:
                texts.append(soup.get_text())
            
            metadata = {
                "source": file_path,
                "file_type": "markdown",
                "headings_count": len(headings),
                "file_metadata": file_metadata,
                "has_tables": bool(soup.find_all('table')),
                "has_lists": bool(soup.find_all(['ul', 'ol']))
            }
            
            return self._create_documents(texts, metadata)
            
        except Exception as e:
            raise Exception(f"Error processing Markdown file {file_path}: {e}")


class DatabaseProcessor(BaseDocumentProcessor):
    """Process SQL databases"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200):
        super().__init__(chunk_size, chunk_overlap)
        self.supported_types = ['sqlite', 'mysql', 'postgresql']
    
    def process(self, connection_string: str, tables: Optional[List[str]] = None) -> List[LangChainDocument]:
        """Extract data from SQL databases"""
        try:
            # For now, supporting SQLite as an example
            if connection_string.endswith('.db') or connection_string.endswith('.sqlite'):
                return self._process_sqlite(connection_string, tables)
            else:
                raise NotImplementedError("Currently only SQLite databases are supported")
            
        except Exception as e:
            raise Exception(f"Error processing database {connection_string}: {e}")
    
    def _process_sqlite(self, db_path: str, tables: Optional[List[str]] = None) -> List[LangChainDocument]:
        """Process SQLite database"""
        conn = sqlite3.connect(db_path)
        cursor = conn.cursor()
        
        try:
            # Get all table names if not specified
            if tables is None:
                cursor.execute("SELECT name FROM sqlite_master WHERE type='table';")
                tables = [row[0] for row in cursor.fetchall()]
            
            texts = []
            tables_info = []
            
            for table_name in tables:
                # Get table schema
                cursor.execute(f"PRAGMA table_info({table_name})")
                columns = cursor.fetchall()
                
                # Get sample data
                cursor.execute(f"SELECT * FROM {table_name} LIMIT 100")
                rows = cursor.fetchall()
                
                # Get row count
                cursor.execute(f"SELECT COUNT(*) FROM {table_name}")
                row_count = cursor.fetchone()[0]
                
                # Create text representation
                table_text = [f"Table: {table_name}"]
                table_text.append(f"Total rows: {row_count}")
                
                # Add schema information
                table_text.append("Schema:")
                for col in columns:
                    table_text.append(f"  {col[1]} ({col[2]})")
                
                # Add sample data
                if rows:
                    table_text.append("\nSample Data:")
                    col_names = [col[1] for col in columns]
                    df = pd.DataFrame(rows, columns=col_names)
                    table_text.append(df.to_string(index=False, max_rows=50))
                
                texts.append("\n".join(table_text))
                
                tables_info.append({
                    "table_name": table_name,
                    "row_count": row_count,
                    "column_count": len(columns),
                    "columns": [{"name": col[1], "type": col[2]} for col in columns]
                })
            
            metadata = {
                "source": db_path,
                "file_type": "database",
                "database_type": "sqlite",
                "total_tables": len(tables),
                "tables_info": tables_info
            }
            
            return self._create_documents(texts, metadata)
            
        finally:
            conn.close()


class EnhancedTextProcessor(BaseDocumentProcessor):
    """Enhanced text processor for .txt files"""
    
    def process(self, file_path: str) -> List[LangChainDocument]:
        """Process text files with encoding detection and structure analysis"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        content = file.read()
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if content is None:
                raise Exception("Could not decode file with any supported encoding")
            
            # Analyze text structure
            lines = content.split('\n')
            paragraphs = [p.strip() for p in content.split('\n\n') if p.strip()]
            
            # Detect if it's structured (like a log file, CSV-like, etc.)
            is_structured = self._detect_structure(lines)
            
            if is_structured:
                texts = self._process_structured_text(lines)
            else:
                texts = paragraphs
            
            metadata = {
                "source": file_path,
                "file_type": "text",
                "encoding": used_encoding,
                "line_count": len(lines),
                "paragraph_count": len(paragraphs),
                "is_structured": is_structured,
                "file_size": os.path.getsize(file_path)
            }
            
            return self._create_documents(texts, metadata)
            
        except Exception as e:
            raise Exception(f"Error processing text file {file_path}: {e}")
    
    def _detect_structure(self, lines: List[str]) -> bool:
        """Detect if text file has a consistent structure"""
        if len(lines) < 3:
            return False
        
        # Check for common delimiters
        delimiters = ['\t', '|', ';', ',']
        for delimiter in delimiters:
            if sum(1 for line in lines[:10] if delimiter in line) > 5:
                return True
        
        # Check for log-like patterns (timestamps, etc.)
        import re
        timestamp_pattern = r'\d{4}-\d{2}-\d{2}|\d{2}:\d{2}:\d{2}'
        timestamp_lines = sum(1 for line in lines[:10] if re.search(timestamp_pattern, line))
        
        return timestamp_lines > 3
    
    def _process_structured_text(self, lines: List[str]) -> List[str]:
        """Process structured text files"""
        # Group related lines together
        chunks = []
        current_chunk = []
        
        for line in lines:
            line = line.strip()
            if not line:
                if current_chunk:
                    chunks.append('\n'.join(current_chunk))
                    current_chunk = []
            else:
                current_chunk.append(line)
                
                # Split on logical boundaries (every 20 lines for structured data)
                if len(current_chunk) >= 20:
                    chunks.append('\n'.join(current_chunk))
                    current_chunk = []
        
        if current_chunk:
            chunks.append('\n'.join(current_chunk))
        
        return chunks


class PDFProcessor(BaseDocumentProcessor):
    """Enhanced PDF processor using PyMuPDF (fitz) with OCR capabilities"""
    
    def process(self, file_path: str) -> List[LangChainDocument]:
        """Process PDF file and extract text with OCR fallback for image-based content"""
        try:
            doc = fitz.open(file_path)
            texts = []
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                
                # First try to extract selectable text
                text = page.get_text()
                page_content = []
                
                # If we have substantial text, use it
                if text.strip() and len(text.strip()) > 50:
                    page_content.append(f"[Text]: {text.strip()}")
                else:
                    # If little or no text, try OCR on page images
                    try:
                        # Convert page to image for OCR
                        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))  # Higher resolution for better OCR
                        img_data = pix.tobytes("png")
                        image = Image.open(io.BytesIO(img_data))
                        
                        # Perform OCR
                        ocr_text = pytesseract.image_to_string(image)
                        if ocr_text.strip():
                            page_content.append(f"[OCR]: {ocr_text.strip()}")
                        
                        # If we still have some selectable text, include it too
                        if text.strip():
                            page_content.append(f"[Text]: {text.strip()}")
                            
                    except Exception as ocr_error:
                        print(f"OCR failed for page {page_num + 1}: {ocr_error}")
                        # Fallback to any available text
                        if text.strip():
                            page_content.append(f"[Text]: {text.strip()}")
                
                # Add page content if we found any
                if page_content:
                    page_text = f"[Page {page_num + 1}]\n" + "\n".join(page_content)
                    texts.append(page_text)
            
            doc.close()
            
            if not texts:
                return []
            
            metadata = {
                'source': file_path,
                'filename': Path(file_path).name,
                'file_type': 'pdf',
                'page_count': len(texts),
                'processor': 'PDFProcessor_OCR'
            }
            
            return self._create_documents(texts, metadata)
            
        except Exception as e:
            print(f"Error processing PDF {file_path}: {e}")
            return []


class DocxProcessor(BaseDocumentProcessor):
    """Enhanced DOCX processor using python-docx"""
    
    def process(self, file_path: str) -> List[LangChainDocument]:
        """Process DOCX file and extract text"""
        try:
            doc = Document(file_path)
            texts = []
            
            # Extract paragraphs
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    texts.append(paragraph.text)
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        texts.append(" | ".join(row_text))
            
            if not texts:
                return []
            
            metadata = {
                'source': file_path,
                'filename': Path(file_path).name,
                'file_type': 'docx',
                'paragraph_count': len([p for p in doc.paragraphs if p.text.strip()]),
                'table_count': len(doc.tables),
                'processor': 'DocxProcessor'
            }
            
            return self._create_documents(texts, metadata)
            
        except Exception as e:
            print(f"Error processing DOCX {file_path}: {e}")
            return []


class DocumentProcessorFactory:
    """Factory class to get the appropriate processor for a file type"""
    
    _processors = {
        '.pdf': PDFProcessor,
        '.docx': DocxProcessor,
        '.pptx': PowerPointProcessor,
        '.ppt': PowerPointProcessor,
        '.xlsx': ExcelProcessor,
        '.csv': ExcelProcessor,
        '.md': MarkdownProcessor,
        '.txt': EnhancedTextProcessor,
        '.db': DatabaseProcessor,
        '.sqlite': DatabaseProcessor
    }
    
    @classmethod
    def get_processor(cls, file_path: str, **kwargs) -> BaseDocumentProcessor:
        """Get the appropriate processor for a file"""
        file_ext = Path(file_path).suffix.lower()
        
        if file_ext not in cls._processors:
            raise ValueError(f"Unsupported file type: {file_ext}")
        
        processor_class = cls._processors[file_ext]
        return processor_class(**kwargs)
    
    @classmethod
    def get_supported_types(cls) -> List[str]:
        """Get list of supported file types"""
        return list(cls._processors.keys())


# Example usage and testing
if __name__ == "__main__":
    # Test the processors
    factory = DocumentProcessorFactory()
    
    print("Supported file types:", factory.get_supported_types())
    
    # Example processing (you would use actual file paths)
    # processor = factory.get_processor("document.pptx")
    # documents = processor.process("document.pptx")
    # print(f"Processed {len(documents)} document chunks")
